import json
import boto3
import io
import urllib.parse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

s3_client = boto3.client('s3')

def lambda_handler(event, context):
    try:
        # Get bucket and object key from event
        bucket_name = event['Records'][0]['s3']['bucket']['name']
        object_key = event['Records'][0]['s3']['object']['key']
        decoded_key = urllib.parse.unquote_plus(object_key)
        print("bucket_name:", bucket_name)
        print("object key:", decoded_key)

        # Get PowerPoint file from S3
        ppt_object = s3_client.get_object(Bucket=bucket_name, Key=decoded_key)
        ppt_content = ppt_object['Body'].read()
        ppt_stream = io.BytesIO(ppt_content)
        
        # Extract text and process images from the PPT stream
        extracted_text = extract_text(ppt_stream)
        process_images(ppt_stream, decoded_key)
        
        # Save extracted text to S3
        save_text(bucket_name, decoded_key, extracted_text)
        
        return {
            'statusCode': 200,
            'body': json.dumps('PPT processed successfully.')
        }

    except Exception as e:
        print(f"Error processing file: {e}")
        return {
            'statusCode': 500,
            'body': json.dumps(f"Error processing file: {str(e)}")
        }

def extract_text(ppt_stream):
    text = ""
    ppt_stream.seek(0)  # Ensure we're at the start of the stream
    presentation = Presentation(ppt_stream)
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += f"Slide {slide_number}: {shape.text}\n"
    return text

def process_images(ppt_stream, decoded_key):
    ppt_stream.seek(0)  # Ensure we're at the start of the stream
    presentation = Presentation(ppt_stream)
    images_key_prefix = decoded_key.replace('.pptx', '')
    
    image_idx = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = io.BytesIO(image.blob)
                
                # Open the image using PIL
                img = Image.open(image_bytes)
                
                # Construct the S3 key for the image
                image_idx += 1
                image_key = f"{images_key_prefix}_image_{image_idx:03d}.{img.format.lower()}"
                
                # Save each image immediately to S3
                s3_client.put_object(Bucket='extractedtextimage', Key=image_key, Body=image_bytes.getvalue(), ContentType=f'image/{img.format.lower()}')

def save_text(bucket_name, decoded_key, extracted_text):
    text_bucket = 'extractedtextimage'
    text_key = decoded_key.replace('.pptx', '.txt')
    
    # Save text to S3 with correct Content-Type
    s3_client.put_object(Bucket=text_bucket, Key=text_key, Body=extracted_text, ContentType='text/plain')
